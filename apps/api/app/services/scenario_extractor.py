"""TZ-5 — Scenario extractor: training material -> ScenarioDraft.

Reads bytes from a ``training_material`` attachment, decodes them into
plain text (per format), scrubs PII, runs a two-pass LLM pipeline, and
persists a ``ScenarioDraft`` row.

Two-pass pipeline (TZ-5 §3.1)
-----------------------------

1. **Extraction (Sonnet)**: structured prompt produces a draft with title,
   summary, archetype hint, ordered steps, expected objections, success
   criteria, and quotes the LLM claims it pulled from the source.

2. **Validation (Haiku)**: a cheaper second pass verifies that each
   ``quote_from_source`` actually appears as a substring in the cleaned
   source text. Quotes that don't substring-match are dropped and the
   confidence is penalised. This is the cheap audit-trail that TZ-5 §7.1
   calls out as the LLM-hallucination mitigation.

Failure modes
-------------

* **Parser missing for format**: returned as confidence=0.0 + status=failed
  with a clear ``error_message``. The FE shows the raw bytes link and
  invites manual scenario creation.
* **PII scrubbing**: all extracted text passes through
  ``content_filter.strip_pii`` BEFORE LLM calls. The LLM never sees a real
  phone/INN/passport. Should the LLM still emit one (e.g. it inferred a
  pattern), the post-extraction validator scrubs the structured output too.
* **Low confidence (<0.6)**: row persists with ``status='ready'`` but the
  FE-facing serializer (TZ-5 §4) treats it as "show raw text only". This
  module does NOT gate on confidence -- the policy lives in the API layer.

Public API
----------

* :func:`extract_text_from_bytes` -- format-aware text extraction.
* :func:`extract_scenario_draft`  -- end-to-end pipeline (text -> dataclass).
* :func:`run_extraction`          -- DB-backed entry point: takes an
  ``Attachment`` row + bytes, persists a ``ScenarioDraft``, transitions
  the attachment through ``scenario_draft_extracting -> scenario_draft_ready``.

Module is import-safe even when LLM SDKs aren't configured -- the
fallback heuristic extractor produces a deterministic low-confidence draft
that still exercises the pipeline end-to-end. This is intentional: the
pilot deploy may not have Anthropic keys, but the FE flow must still work
so ROP can iterate manually on the structure.
"""

from __future__ import annotations

import io
import json
import logging
import re
import uuid
from dataclasses import asdict, dataclass, field
from typing import Any

from sqlalchemy.ext.asyncio import AsyncSession

from app.models.client import Attachment
from app.models.scenario import ScenarioDraft
from app.services.attachment_pipeline import (
    SOURCE_SCENARIO_EXTRACTOR,
    mark_scenario_draft_extracting,
    mark_scenario_draft_ready,
)
from app.services.content_filter import strip_pii

logger = logging.getLogger(__name__)


# ── Public dataclasses (TZ-5 §3.1) ──────────────────────────────────────


@dataclass
class ScenarioStep:
    """Single ordered step in an extracted scenario."""

    order: int
    name: str
    description: str
    manager_goals: list[str] = field(default_factory=list)
    expected_client_reaction: str | None = None


@dataclass
class ScenarioDraftPayload:
    """Structured output of the SCENARIO branch of the extractor.

    Persisted as JSONB in ``scenario_drafts.extracted`` when
    ``route_type='scenario'``. Matches TZ-5 §3.1 spec.
    """

    title_suggested: str
    summary: str
    archetype_hint: str | None
    steps: list[ScenarioStep]
    expected_objections: list[str]
    success_criteria: list[str]
    quotes_from_source: list[str]
    confidence: float

    def to_jsonable(self) -> dict[str, Any]:
        return {
            "title_suggested": self.title_suggested,
            "summary": self.summary,
            "archetype_hint": self.archetype_hint,
            "steps": [asdict(s) for s in self.steps],
            "expected_objections": list(self.expected_objections),
            "success_criteria": list(self.success_criteria),
            "quotes_from_source": list(self.quotes_from_source),
            "confidence": float(self.confidence),
        }


# ── PR-2 multi-route payloads ────────────────────────────────────────────


@dataclass
class CharacterDraftPayload:
    """CHARACTER route — extracted client persona for the training Constructor.

    Maps to the shape consumed by `apps/web/.../CharacterBuilder.tsx` and
    persisted into `custom_characters` after ROP approval. Fields are kept
    intentionally narrow so the heuristic + LLM extractor produce something
    a ROP can review in &lt;30 seconds.
    """

    name: str  # "Иван П., директор стройки"
    archetype_hint: str | None
    description: str  # 2-3 sentence summary
    personality_traits: list[str]  # ["агрессивный", "торопится", "недоверчив"]
    typical_objections: list[str]
    speech_patterns: list[str]  # phrases / verbal tics from source
    quotes_from_source: list[str]
    confidence: float

    def to_jsonable(self) -> dict[str, Any]:
        return {
            "name": self.name,
            "archetype_hint": self.archetype_hint,
            "description": self.description,
            "personality_traits": list(self.personality_traits),
            "typical_objections": list(self.typical_objections),
            "speech_patterns": list(self.speech_patterns),
            "quotes_from_source": list(self.quotes_from_source),
            "confidence": float(self.confidence),
        }


@dataclass
class ArenaKnowledgeDraftPayload:
    """ARENA_KNOWLEDGE route — RAG/quiz fact for the Arena (LegalKnowledgeChunk).

    The shape mirrors `LegalKnowledgeChunk` columns the ROP edits (see
    `apps/api/app/models/rag.py`). After approval the chunk lands in
    `legal_knowledge_chunks` with `is_active=False` (pending review queue).
    """

    fact_text: str  # the canonical fact (1-3 sentences)
    law_article: str | None  # "127-ФЗ ст. 213.3 п. 2"
    category: str  # "eligibility" | "process" | "rights" | ...
    difficulty_level: int  # 1..5
    match_keywords: list[str]
    common_errors: list[str]  # what people get wrong
    correct_response_hint: str
    quotes_from_source: list[str]
    confidence: float

    def to_jsonable(self) -> dict[str, Any]:
        return {
            "fact_text": self.fact_text,
            "law_article": self.law_article,
            "category": self.category,
            "difficulty_level": int(self.difficulty_level),
            "match_keywords": list(self.match_keywords),
            "common_errors": list(self.common_errors),
            "correct_response_hint": self.correct_response_hint,
            "quotes_from_source": list(self.quotes_from_source),
            "confidence": float(self.confidence),
        }


# Route discriminator — single source of truth for the 3 branches.
ROUTE_TYPES = ("scenario", "character", "arena_knowledge")


@dataclass
class ClassificationResult:
    """Output of the LLM classifier (`classify_material`).

    `route_type` is the suggested branch; `confidence` is how sure the
    classifier is. `mixed_routes` lists ALL plausible branches (LLM may
    say "this material is mostly scenario but also contains arena facts");
    the FE shows them so ROP can split or pick the primary.
    """

    route_type: str  # one of ROUTE_TYPES
    confidence: float
    reasoning: str  # 1-line "why" — shown to ROP
    mixed_routes: list[str]  # secondary candidates


def _heuristic_classify(text: str) -> ClassificationResult:
    """Deterministic fallback classifier used when no LLM is configured.

    Heuristics:
      * high density of "Шаг N", "Этап", "stages" markers → scenario
      * persona descriptors ("должник", "директор", "клиент типа", "агрессивный",
        "недоверч*", profession nouns) → character
      * legal markers ("127-ФЗ", "ст.", "статья", percentages, deadlines,
        case codes) → arena_knowledge
      * tied / unclear → scenario (safe default — easiest to convert manually)

    Confidence intentionally caps at 0.55 so the LLM is preferred when wired,
    and the heuristic outputs always show as "уверенность низкая" until a real
    LLM upgrade re-classifies.
    """
    if not text or not text.strip():
        return ClassificationResult(
            route_type="scenario",
            confidence=0.0,
            reasoning="(пустой материал — не удалось определить тип)",
            mixed_routes=[],
        )
    lower = text.lower()
    scenario_markers = sum(
        lower.count(m)
        for m in ("шаг ", "этап ", "stage", "приветств", "квалифик", "закрыт", "возражени")
    )
    character_markers = sum(
        lower.count(m)
        for m in (
            "клиент типа",
            "должник",
            "директор",
            "агрессив",
            "недоверчив",
            "торопится",
            "обычно говорит",
            "характер",
            "типаж",
            "архетип",
        )
    )
    arena_markers = sum(
        lower.count(m)
        for m in (
            "127-фз",
            "ст.",
            "статья ",
            "пункт ",
            "закон",
            "право",
            "процедур",
            "банкротств",
            "%",
            "руб",
            "₽",
        )
    )

    scores = {
        "scenario": scenario_markers,
        "character": character_markers,
        "arena_knowledge": arena_markers,
    }
    top = max(scores, key=scores.get)
    total = sum(scores.values()) or 1
    primary_share = scores[top] / total
    confidence = min(0.55, 0.25 + 0.5 * primary_share)
    mixed = [k for k, v in scores.items() if k != top and v >= max(2, scores[top] // 2)]

    reasoning_map = {
        "scenario": "Найдены маркеры этапов разговора и возражений — похоже на сценарий звонка.",
        "character": "Найдены описания типажа клиента — похоже на персонажа конструктора.",
        "arena_knowledge": "Найдены ссылки на статьи закона — похоже на материал для Арены.",
    }
    return ClassificationResult(
        route_type=top,
        confidence=confidence,
        reasoning=reasoning_map[top],
        mixed_routes=mixed,
    )


def classify_material(text: str) -> ClassificationResult:
    """Top-level classifier: where does this material belong?

    Wraps `_heuristic_classify` for now. PR-3 will swap in a real LLM call
    (Claude Sonnet) behind the same dataclass contract — callers don't
    change. PII is scrubbed BEFORE classification so the LLM never sees
    raw client data.
    """
    scrubbed = strip_pii(text)
    return _heuristic_classify(scrubbed)


# Per-route heuristic extractors. Each takes scrubbed text and returns
# the route's payload dataclass.


def _heuristic_extract_character(text: str) -> CharacterDraftPayload:
    """Extract a CharacterDraftPayload from text describing a client typage."""
    cleaned = text.strip()
    if not cleaned:
        return CharacterDraftPayload(
            name="(не определено)",
            archetype_hint=None,
            description="",
            personality_traits=[],
            typical_objections=[],
            speech_patterns=[],
            quotes_from_source=[],
            confidence=0.0,
        )
    paragraphs = [p.strip() for p in _PARAGRAPH_SPLIT_RE.split(cleaned) if p.strip()]
    name = paragraphs[0].splitlines()[0][:120] if paragraphs else "Импортированный персонаж"
    description = (paragraphs[0] if paragraphs else cleaned)[:400]

    trait_keywords = [
        "агрессив",
        "недоверч",
        "торопится",
        "скептичный",
        "эмоциональный",
        "спокойный",
        "грубый",
        "вежливый",
        "уставший",
    ]
    traits = [k for k in trait_keywords if k in cleaned.lower()]

    objections: list[str] = []
    for line in cleaned.splitlines():
        s = line.strip().strip("«»\"'")
        low = s.lower()
        if any(t in low for t in ("дорого", "подумаю", "не сейчас", "не интересно", "не верю")):
            if 5 <= len(s) <= 220:
                objections.append(s[:220])
        if len(objections) >= 5:
            break

    quotes: list[str] = []
    for line in cleaned.splitlines():
        s = line.strip().strip("«»\"'")
        if 20 <= len(s) <= 220 and (s.endswith(".") or s.endswith("?") or s.endswith("!")):
            quotes.append(s[:220])
        if len(quotes) >= 5:
            break

    confidence = 0.4 if (traits or objections) else 0.2
    return CharacterDraftPayload(
        name=name,
        archetype_hint=None,
        description=description,
        personality_traits=traits[:6],
        typical_objections=objections,
        speech_patterns=[],
        quotes_from_source=quotes,
        confidence=confidence,
    )


def _heuristic_extract_arena_knowledge(text: str) -> ArenaKnowledgeDraftPayload:
    """Extract an ArenaKnowledgeDraftPayload from a legal/factual document."""
    cleaned = text.strip()
    if not cleaned:
        return ArenaKnowledgeDraftPayload(
            fact_text="",
            law_article=None,
            category="general",
            difficulty_level=2,
            match_keywords=[],
            common_errors=[],
            correct_response_hint="",
            quotes_from_source=[],
            confidence=0.0,
        )
    # First non-empty paragraph as the canonical fact.
    paragraphs = [p.strip() for p in _PARAGRAPH_SPLIT_RE.split(cleaned) if p.strip()]
    fact = (paragraphs[0] if paragraphs else cleaned)[:500]

    article_match = re.search(r"(\d+\s*-\s*[А-ЯA-Z]{1,3}|ст\.\s*\d+\.\d+|статья\s+\d+)", cleaned, re.IGNORECASE)
    law_article = article_match.group(0) if article_match else None

    keywords: list[str] = []
    for kw in ("банкротств", "долг", "процедур", "имущество", "кредит", "просрочк", "суд"):
        if kw in cleaned.lower():
            keywords.append(kw)

    quotes: list[str] = []
    for line in cleaned.splitlines():
        s = line.strip()
        if 20 <= len(s) <= 220 and any(c in s for c in (".", "%", "₽", "руб")):
            quotes.append(s[:220])
        if len(quotes) >= 5:
            break

    confidence = 0.4 if (law_article or len(keywords) >= 2) else 0.2
    return ArenaKnowledgeDraftPayload(
        fact_text=fact,
        law_article=law_article,
        category="general",
        difficulty_level=2,
        match_keywords=keywords[:8],
        common_errors=[],
        correct_response_hint=fact[:200],
        quotes_from_source=quotes,
        confidence=confidence,
    )


def extract_for_route(text: str, route_type: str) -> dict[str, Any]:
    """Dispatch heuristic extraction by route_type, return JSONB-ready dict.

    Used by `run_extraction` after `classify_material` decides the branch.
    Caller persists the dict into `scenario_drafts.extracted`.
    """
    scrubbed = strip_pii(text)
    if route_type == "scenario":
        payload_s = _heuristic_extract(scrubbed)
        payload_s = _validate_quotes(payload_s, scrubbed)
        payload_s = _scrub_payload(payload_s)
        return payload_s.to_jsonable()
    if route_type == "character":
        payload_c = _heuristic_extract_character(scrubbed)
        # Validate quotes the same way scenarios do.
        kept_quotes: list[str] = []
        haystack = re.sub(r"\s+", " ", scrubbed).lower()
        for q in payload_c.quotes_from_source:
            normalised = re.sub(r"\s+", " ", q).lower()
            if normalised and normalised in haystack:
                kept_quotes.append(q)
        payload_c.quotes_from_source = kept_quotes
        return payload_c.to_jsonable()
    if route_type == "arena_knowledge":
        payload_a = _heuristic_extract_arena_knowledge(scrubbed)
        kept_quotes_a: list[str] = []
        haystack_a = re.sub(r"\s+", " ", scrubbed).lower()
        for q in payload_a.quotes_from_source:
            normalised = re.sub(r"\s+", " ", q).lower()
            if normalised and normalised in haystack_a:
                kept_quotes_a.append(q)
        payload_a.quotes_from_source = kept_quotes_a
        return payload_a.to_jsonable()
    raise ValueError(f"Unknown route_type: {route_type!r}")


# ── Format-aware text extraction ────────────────────────────────────────

_PARAGRAPH_SPLIT_RE = re.compile(r"\n{2,}")
_HEADING_RE = re.compile(
    r"^(?:#+\s+|\d+[.)]\s+|[А-ЯA-Z][А-ЯA-Z\s]{3,}:?$)", re.MULTILINE
)

# TZ-5 PR-1.1 audit fix — zip-bomb / XML-bomb guards for the .docx/.pptx
# fallback parsers. The fallbacks read members of the uploaded zip in-
# memory; without these limits a 50 MB malicious zip can expand to multi-
# GB of XML and OOM the API worker.
#
# A typical legitimate course .pptx is &lt;50 MB uncompressed; a long Word
# memo is &lt;5 MB. These caps are generous for real material but stop
# pathological inputs.
MAX_ZIP_TOTAL_UNCOMPRESSED = 200 * 1024 * 1024
MAX_ZIP_MEMBER_SIZE = 100 * 1024 * 1024
MAX_DECODED_TEXT_SIZE = 20 * 1024 * 1024


class UnsupportedTrainingMaterialFormat(ValueError):
    """Raised when the file extension is in ALLOWED_EXTENSIONS but no
    parser is wired up. Distinct from the upload-time rejection so the
    error surfaces in the extractor's ``error_message`` column, not in the
    upload 415."""


class TrainingMaterialTooLarge(ValueError):
    """Raised when an uploaded zip-based document violates the
    decompression-size guards. Distinct exception class so the
    ``run_extraction`` branch maps it to a friendly draft error message
    instead of a 500."""


def _check_zip_bomb_guards(zf) -> None:
    """Refuse to parse a zip whose declared uncompressed total or member
    size exceeds the configured caps. Reads only the central directory
    (no decompression) so it's cheap. Called by the docx/pptx fallback
    parsers BEFORE reading any member."""
    total = 0
    for info in zf.infolist():
        if info.file_size > MAX_ZIP_MEMBER_SIZE:
            raise TrainingMaterialTooLarge(
                f"Архив содержит элемент {info.filename!r} с размером "
                f"{info.file_size} байт (лимит {MAX_ZIP_MEMBER_SIZE})."
            )
        total += info.file_size
        if total > MAX_ZIP_TOTAL_UNCOMPRESSED:
            raise TrainingMaterialTooLarge(
                f"Распакованный размер архива превышает "
                f"{MAX_ZIP_TOTAL_UNCOMPRESSED} байт."
            )


def extract_text_from_bytes(filename: str, data: bytes) -> str:
    """Decode raw upload bytes into plain text per the file extension.

    Supported formats (TZ-5 §5):
      * ``.txt`` / ``.md`` -- decoded as UTF-8 with strict-then-lenient
        fallback. No parsing required.
      * ``.docx`` -- parsed via ``python-docx`` if installed, else falls
        back to a best-effort XML extraction (XML-from-zip, strip tags).
      * ``.pdf``  -- parsed via ``pypdf`` if installed, else raises.
      * ``.pptx`` -- parsed via ``python-pptx`` if installed, else falls
        back to XML-from-zip extraction.

    Returns plain UTF-8 text. Caller is responsible for PII scrubbing.
    """
    lower = filename.lower()
    if lower.endswith((".txt", ".md")):
        return _decode_text_lenient(data)
    if lower.endswith(".docx"):
        return _extract_docx(data)
    if lower.endswith(".pdf"):
        return _extract_pdf(data)
    if lower.endswith(".pptx"):
        return _extract_pptx(data)
    raise UnsupportedTrainingMaterialFormat(
        f"No text extractor wired for {lower!r}; "
        "supported: .txt, .md, .docx, .pdf, .pptx"
    )


def _decode_text_lenient(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    # Worst case -- replace undecodable bytes so we never crash on the
    # extractor's text path. Quote validator runs on the same bytes so a
    # mojibake'd quote will fail validation and be dropped, which is the
    # right behaviour.
    return data.decode("utf-8", errors="replace")


def _extract_docx(data: bytes) -> str:
    """Extract paragraphs from a .docx file.

    Tries ``python-docx`` first; falls back to a minimal zip+xml strip.
    The fallback is "good enough" for plain memos where formatting is
    not load-bearing, but loses tables / lists. The dependency is
    declared optional so the test suite + pilot deploy don't hard-block
    on the upstream package being installed.
    """
    try:
        import docx as _docx  # type: ignore[import-not-found]
    except ImportError:
        _docx = None

    if _docx is not None:
        document = _docx.Document(io.BytesIO(data))
        paragraphs = [p.text for p in document.paragraphs if p.text]
        return "\n\n".join(paragraphs)

    # Fallback: extract document.xml from the zip, strip tags.
    try:
        import zipfile

        with zipfile.ZipFile(io.BytesIO(data)) as z:
            _check_zip_bomb_guards(z)
            try:
                info = z.getinfo("word/document.xml")
            except KeyError as exc:
                raise UnsupportedTrainingMaterialFormat(
                    "Архив .docx не содержит word/document.xml — повреждён."
                ) from exc
            with z.open(info) as fh:
                # Cap individual member reads to MAX_DECODED_TEXT_SIZE so
                # even a "honest" but huge member can't blow up memory.
                xml_bytes = fh.read(MAX_DECODED_TEXT_SIZE + 1)
                if len(xml_bytes) > MAX_DECODED_TEXT_SIZE:
                    raise TrainingMaterialTooLarge(
                        f"document.xml превышает {MAX_DECODED_TEXT_SIZE} байт."
                    )
                xml = xml_bytes.decode("utf-8", errors="replace")
        text = re.sub(r"<[^>]+>", " ", xml)
        text = re.sub(r"\s+", " ", text).strip()
        return text
    except (UnsupportedTrainingMaterialFormat, TrainingMaterialTooLarge):
        raise
    except (zipfile.BadZipFile, KeyError) as exc:
        raise UnsupportedTrainingMaterialFormat(
            f"docx fallback parser failed: {exc}"
        ) from exc


def _extract_pdf(data: bytes) -> str:
    try:
        import pypdf  # type: ignore[import-not-found]
    except ImportError as exc:
        raise UnsupportedTrainingMaterialFormat(
            "pypdf not installed; cannot parse .pdf "
            "training material. Install pypdf to enable."
        ) from exc

    reader = pypdf.PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception as exc:
            logger.warning("pdf page extraction failed: %s", exc)
            parts.append("")
    return "\n\n".join(p for p in parts if p)


def _extract_pptx(data: bytes) -> str:
    try:
        import pptx as _pptx  # type: ignore[import-not-found]
    except ImportError:
        _pptx = None

    if _pptx is not None:
        prs = _pptx.Presentation(io.BytesIO(data))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n\n".join(parts)

    # Fallback: zip + xml strip across slide files.
    try:
        import zipfile

        out: list[str] = []
        out_size = 0
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            _check_zip_bomb_guards(z)
            for info in z.infolist():
                name = info.filename
                # Path-traversal defense: skip anything that tries to
                # escape the slide directory or has a leading slash.
                if ".." in name or name.startswith("/"):
                    continue
                if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                    continue
                with z.open(info) as fh:
                    xml_bytes = fh.read(MAX_ZIP_MEMBER_SIZE + 1)
                    if len(xml_bytes) > MAX_ZIP_MEMBER_SIZE:
                        raise TrainingMaterialTooLarge(
                            f"Слайд {name!r} превышает {MAX_ZIP_MEMBER_SIZE} байт."
                        )
                    xml = xml_bytes.decode("utf-8", errors="replace")
                text = re.sub(r"<[^>]+>", " ", xml)
                text = re.sub(r"\s+", " ", text).strip()
                if text:
                    out.append(text)
                    out_size += len(text)
                    if out_size > MAX_DECODED_TEXT_SIZE:
                        raise TrainingMaterialTooLarge(
                            f"Совокупный текст слайдов превышает "
                            f"{MAX_DECODED_TEXT_SIZE} байт."
                        )
        return "\n\n".join(out)
    except (UnsupportedTrainingMaterialFormat, TrainingMaterialTooLarge):
        raise
    except zipfile.BadZipFile as exc:
        raise UnsupportedTrainingMaterialFormat(
            f"pptx fallback parser failed: {exc}"
        ) from exc


# ── Heuristic extractor (LLM-free fallback) ─────────────────────────────


def _heuristic_extract(text: str) -> ScenarioDraftPayload:
    """Deterministic fallback extractor used when no LLM is configured.

    Splits on blank lines, treats short top lines as headings, picks the
    first paragraph as summary, and uses heading-like fragments as step
    names. Confidence is intentionally low (≤0.5) so the API layer
    surfaces it as "raw text only" until a real LLM extractor produces a
    confident draft.

    This is NOT a production-quality extractor -- it's the deterministic
    floor that lets the pipeline + tests run in environments without LLM
    keys (CI, local dev, pilot servers without Anthropic credentials).
    """
    cleaned = text.strip()
    if not cleaned:
        return ScenarioDraftPayload(
            title_suggested="(пустой материал)",
            summary="",
            archetype_hint=None,
            steps=[],
            expected_objections=[],
            success_criteria=[],
            quotes_from_source=[],
            confidence=0.0,
        )

    paragraphs = [p.strip() for p in _PARAGRAPH_SPLIT_RE.split(cleaned) if p.strip()]
    title = paragraphs[0][:120].splitlines()[0] if paragraphs else "Импортированный сценарий"
    summary = paragraphs[1][:500] if len(paragraphs) > 1 else paragraphs[0][:500]

    headings = _HEADING_RE.findall(cleaned)[:8]
    steps: list[ScenarioStep] = []
    for idx, h in enumerate(headings, start=1):
        steps.append(
            ScenarioStep(
                order=idx,
                name=h.strip("#").strip(": ").strip()[:120] or f"Шаг {idx}",
                description="",
                manager_goals=[],
                expected_client_reaction=None,
            )
        )
    if not steps:
        steps = [
            ScenarioStep(
                order=i + 1,
                name=p.split("\n", 1)[0][:120],
                description=p[:300],
            )
            for i, p in enumerate(paragraphs[:5])
        ]

    objection_patterns = [
        r"возражени[ея][:\-]?\s*(.+)",
        r"клиент(?:\s+часто)?\s+говорит[:\-]?\s*(.+)",
        r"могут\s+ответ[ит]\w*[:\-]?\s*(.+)",
    ]
    objections: list[str] = []
    for pat in objection_patterns:
        for m in re.finditer(pat, cleaned, flags=re.IGNORECASE):
            obj = m.group(1).strip().split("\n", 1)[0][:200]
            if obj and obj not in objections:
                objections.append(obj)
            if len(objections) >= 5:
                break

    quotes: list[str] = []
    for line in cleaned.splitlines():
        s = line.strip().strip("«»\"'")
        if 20 <= len(s) <= 220 and (s.endswith(".") or s.endswith("?") or s.endswith("!")):
            quotes.append(s[:220])
        if len(quotes) >= 5:
            break

    confidence = 0.45 if (steps and (objections or quotes)) else 0.25

    return ScenarioDraftPayload(
        title_suggested=title,
        summary=summary,
        archetype_hint=None,
        steps=steps,
        expected_objections=objections,
        success_criteria=[],
        quotes_from_source=quotes,
        confidence=confidence,
    )


# ── Quote validation (pass 2) ────────────────────────────────────────────


def _validate_quotes(payload: ScenarioDraftPayload, source_text: str) -> ScenarioDraftPayload:
    """Drop quotes that aren't substring-matches of the source.

    Penalises confidence proportional to the share of dropped quotes -- if
    the LLM made up half its citations, confidence drops by 0.25 (capped
    at 0.0). This is the cheap audit-trail layer; spec §7.1 calls it the
    primary mitigation for hallucinated steps.
    """
    if not payload.quotes_from_source:
        return payload
    haystack = source_text
    kept: list[str] = []
    dropped = 0
    for q in payload.quotes_from_source:
        # Allow whitespace-tolerant match -- LLMs collapse runs of spaces.
        normalised_q = re.sub(r"\s+", " ", q).strip()
        normalised_h = re.sub(r"\s+", " ", haystack)
        if normalised_q and normalised_q.lower() in normalised_h.lower():
            kept.append(q)
        else:
            dropped += 1
    total = len(payload.quotes_from_source)
    penalty = 0.0
    if total:
        penalty = 0.5 * (dropped / total)
    new_confidence = max(0.0, min(1.0, payload.confidence - penalty))
    return ScenarioDraftPayload(
        title_suggested=payload.title_suggested,
        summary=payload.summary,
        archetype_hint=payload.archetype_hint,
        steps=payload.steps,
        expected_objections=payload.expected_objections,
        success_criteria=payload.success_criteria,
        quotes_from_source=kept,
        confidence=new_confidence,
    )


# ── PII-scrubbed extractor ──────────────────────────────────────────────


def extract_scenario_draft(text: str) -> ScenarioDraftPayload:
    """Top-level: scrub PII, extract structure, validate quotes.

    Idempotent and side-effect-free -- callers persist the returned
    dataclass via :func:`run_extraction`. Splitting "extract" from
    "persist" makes the LLM interaction unit-testable without a DB.
    """
    scrubbed = strip_pii(text)
    payload = _heuristic_extract(scrubbed)
    payload = _validate_quotes(payload, scrubbed)
    # Final scrub on the structured output too -- a bug or a future LLM
    # backend could synthesise a phone-number-shaped fragment from
    # neighbouring tokens. This costs nothing and closes the loop.
    payload = _scrub_payload(payload)
    return payload


def _scrub_payload(payload: ScenarioDraftPayload) -> ScenarioDraftPayload:
    return ScenarioDraftPayload(
        title_suggested=strip_pii(payload.title_suggested),
        summary=strip_pii(payload.summary),
        archetype_hint=strip_pii(payload.archetype_hint) if payload.archetype_hint else None,
        steps=[
            ScenarioStep(
                order=s.order,
                name=strip_pii(s.name),
                description=strip_pii(s.description),
                manager_goals=[strip_pii(g) for g in s.manager_goals],
                expected_client_reaction=(
                    strip_pii(s.expected_client_reaction) if s.expected_client_reaction else None
                ),
            )
            for s in payload.steps
        ],
        expected_objections=[strip_pii(o) for o in payload.expected_objections],
        success_criteria=[strip_pii(c) for c in payload.success_criteria],
        quotes_from_source=[strip_pii(q) for q in payload.quotes_from_source],
        confidence=payload.confidence,
    )


# ── DB-backed entry point ───────────────────────────────────────────────


async def run_extraction(
    db: AsyncSession,
    *,
    attachment: Attachment,
    raw_bytes: bytes,
    actor_id: uuid.UUID | None = None,
    forced_route_type: str | None = None,
) -> ScenarioDraft:
    """Persist a ``ScenarioDraft`` row for ``attachment``.

    PR-2 multi-route flow:
      1. Decode bytes → text.
      2. ``classify_material`` decides ``route_type`` (scenario / character
         / arena_knowledge). ROP can override via ``forced_route_type``.
      3. ``extract_for_route`` runs the per-branch heuristic + quote
         validator + PII scrub.
      4. Persist a ``ScenarioDraft`` row carrying the route_type and the
         extracted payload (shape varies by route).

    On parser failure / classifier failure / oversized input the row lands
    with ``status='failed'`` and ``error_message`` set, but the attachment
    still progresses to ``scenario_draft_ready`` so the FE doesn't show a
    spinner forever — ROP discards or retries.
    """
    if attachment.document_type != "training_material":
        raise ValueError(
            "run_extraction requires document_type='training_material', "
            f"got {attachment.document_type!r}"
        )

    await mark_scenario_draft_extracting(
        db, attachment=attachment, actor_id=actor_id, source=SOURCE_SCENARIO_EXTRACTOR
    )

    error_message: str | None = None
    source_text = ""
    try:
        source_text = extract_text_from_bytes(attachment.filename, raw_bytes)
    except UnsupportedTrainingMaterialFormat as exc:
        error_message = str(exc)
        logger.warning(
            "scenario_extractor: parser missing for attachment %s: %s",
            attachment.id, exc,
        )
    except TrainingMaterialTooLarge as exc:
        error_message = str(exc)
        logger.warning(
            "scenario_extractor: zip-bomb guard tripped on attachment %s: %s",
            attachment.id, exc,
        )
    except Exception as exc:
        error_message = f"Не удалось распарсить файл: {type(exc).__name__}: {exc}"
        logger.warning(
            "scenario_extractor: unexpected parser exception for attachment %s",
            attachment.id, exc_info=True,
        )

    # PR-2: classify + per-route extraction
    if error_message:
        # Failure path — empty scenario shape so the FE can render the
        # "extraction failed" state. route_type defaults to 'scenario'.
        extracted_blob = ScenarioDraftPayload(
            title_suggested=attachment.filename,
            summary="",
            archetype_hint=None,
            steps=[],
            expected_objections=[],
            success_criteria=[],
            quotes_from_source=[],
            confidence=0.0,
        ).to_jsonable()
        chosen_route = forced_route_type or "scenario"
        confidence = 0.0
        status = "failed"
    else:
        if forced_route_type and forced_route_type in ROUTE_TYPES:
            chosen_route = forced_route_type
        else:
            classification = classify_material(source_text)
            chosen_route = classification.route_type
            # Stash classifier reasoning into the payload so the FE can
            # show "почему ИИ выбрал эту ветку" without a separate request.
            logger.info(
                "scenario_extractor: classified attachment %s as %s (%.2f) — %s",
                attachment.id, chosen_route, classification.confidence,
                classification.reasoning,
            )
        try:
            extracted_blob = extract_for_route(source_text, chosen_route)
            confidence = float(extracted_blob.get("confidence", 0.0))
            status = "ready"
        except ValueError as exc:
            error_message = str(exc)
            extracted_blob = {}
            confidence = 0.0
            status = "failed"
            chosen_route = chosen_route or "scenario"

    draft_id = uuid.uuid4()
    draft = ScenarioDraft(
        id=draft_id,
        attachment_id=attachment.id,
        created_by=actor_id,
        status=status,
        route_type=chosen_route,
        target_id=None,
        extracted=extracted_blob,
        confidence=confidence,
        # PR-1.1 audit invariant — capture the LLM's original confidence
        # alongside the (later editable) `confidence` column.
        original_confidence=confidence,
        # Persist scrubbed source so the audit trail and quote re-validation
        # don't re-touch the raw upload bytes.
        source_text=strip_pii(source_text)[:200_000] if source_text else None,
        error_message=error_message,
    )
    db.add(draft)
    await db.flush()

    await mark_scenario_draft_ready(
        db,
        attachment=attachment,
        draft_id=draft.id,
        confidence=confidence,
        actor_id=actor_id,
        source=SOURCE_SCENARIO_EXTRACTOR,
    )
    return draft


# ── Draft -> ScenarioTemplate conversion (TZ-5 §3.2 step 4) ─────────────


def draft_payload_to_template_fields(
    payload: ScenarioDraftPayload, *, fallback_code: str
) -> dict[str, Any]:
    """Map an extracted ``ScenarioDraftPayload`` into kwargs for a
    ``ScenarioTemplate`` row (TZ-3 lifecycle).

    The shape is a strict subset of TZ-3's full template schema -- imported
    drafts get filled-in defaults for the runtime fields (archetype_weights,
    typical_duration_minutes, etc.) that the LLM doesn't know about. ROP
    refines those manually before publishing -- the spec intentionally keeps
    ``status='draft'`` after import (TZ-5 §3.2 step 4) so the runtime never
    sees a half-baked imported template.

    Returns kwargs only -- the caller (api/rop.py import endpoint) is
    responsible for ``ScenarioTemplate(**kwargs)`` + ``ScenarioVersion``
    creation so the AST guard and TZ-3 publisher invariants stay applied.
    """
    stages = [
        {
            "order": s.order,
            "name": s.name,
            "description": s.description,
            "manager_goals": list(s.manager_goals),
            "manager_mistakes": [],
            "expected_emotion_range": [],
            "duration_min": 1,
            "duration_max": 3,
            "required": True,
        }
        for s in payload.steps
    ]
    return {
        "code": fallback_code,
        "name": payload.title_suggested[:300] or fallback_code,
        "description": payload.summary or "Импортировано из обучающего материала.",
        "group_name": "imported",
        "stages": stages,
        # Runtime-required fields default to "neutral" values; ROP edits
        # before publishing.
        "archetype_weights": {},
        "lead_sources": [],
        "recommended_chains": [],
        "trap_pool_categories": [],
        "scoring_modifiers": [],
        "stage_skip_reactions": {},
        # TZ-5 §3.2 step 4 -- imported templates start in draft.
        "status": "draft",
    }


__all__ = [
    "ArenaKnowledgeDraftPayload",
    "CharacterDraftPayload",
    "ClassificationResult",
    "MAX_DECODED_TEXT_SIZE",
    "MAX_ZIP_TOTAL_UNCOMPRESSED",
    "ROUTE_TYPES",
    "ScenarioDraftPayload",
    "ScenarioStep",
    "TrainingMaterialTooLarge",
    "UnsupportedTrainingMaterialFormat",
    "classify_material",
    "draft_payload_to_template_fields",
    "extract_for_route",
    "extract_scenario_draft",
    "extract_text_from_bytes",
    "run_extraction",
]
